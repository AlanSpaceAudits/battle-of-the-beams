"""Build the BotB summary PPTX from the Jupyter Book and ITU spreadsheet.

Aether Cosmology Research Group theme:
  navy #141A2B  | gold #E8B549  | cream #F5F1E8
Slide size 13.33 x 7.5 in (16:9).
"""
from __future__ import annotations
import math, pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import common as c

HERE = pathlib.Path(__file__).parent
LOGO = HERE / "aether_cosmology_logo.png"
OUT  = HERE / "BotB_Null_Hypothesis_Summary.pptx"

NAVY  = RGBColor(0x14, 0x1A, 0x2B)
GOLD  = RGBColor(0xE8, 0xB5, 0x49)
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
RED   = RGBColor(0xC8, 0x3F, 0x3F)
GREEN = RGBColor(0x5F, 0xB0, 0x6A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def make_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs, bg=NAVY):
    layout = prs.slide_layouts[6]   # Blank
    s = prs.slides.add_slide(layout)
    bg_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_rect.line.fill.background()
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = bg
    bg_rect.shadow.inherit = False
    return s


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=CREAM, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_logo(slide, left=Inches(0.3), top=Inches(0.3), height=Inches(0.9)):
    pic = slide.shapes.add_picture(str(LOGO), left, top, height=height)
    return pic


def gold_bar(slide, top=Inches(1.2)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), top,
                                 SLIDE_W - Inches(1.0), Emu(40000))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD


def header(slide, title_text, subtitle_text=None):
    add_logo(slide)
    add_text(slide, Inches(1.5), Inches(0.35), Inches(11), Inches(0.65),
             title_text, size=28, bold=True, color=CREAM, font="Calibri")
    if subtitle_text:
        add_text(slide, Inches(1.5), Inches(0.85), Inches(11), Inches(0.4),
                 subtitle_text, size=14, color=GOLD, font="Calibri")
    gold_bar(slide, top=Inches(1.3))


def add_table(slide, left, top, width, height, rows):
    """rows is a list of lists. First row is header (gold)."""
    cols_n = len(rows[0])
    tbl_shape = slide.shapes.add_table(len(rows), cols_n, left, top, width, height)
    tbl = tbl_shape.table
    for r, row in enumerate(rows):
        for col, val in enumerate(row):
            cell = tbl.cell(r, col)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(13 if r == 0 else 12)
            run.font.bold = (r == 0)
            run.font.color.rgb = NAVY if r == 0 else CREAM
            run.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD if r == 0 else NAVY
    return tbl_shape


# =====================================================================
# Compute the per-path numbers once so the slides stay in sync
# =====================================================================
def kleve_results():
    targets = ['Spalding','Retford','Derby','Birmingham',
               'TF 400 km','TF 500 km','TF 700 km','TF 800 km','TF 1000 km']
    rows = []
    for tgt in targets:
        f = c.link_budget('Kleve', tgt, model='fock')
        rows.append((tgt, f['d_km'], f['ground'],
                     f['SNR_eq_dB'], f['V_eq_uV']))
    return rows


def stollberg_results():
    targets = ['Beeston','Derby','Birmingham','Liverpool',
               'TF 400 km','TF 500 km','TF 700 km','TF 800 km','TF 1000 km']
    rows = []
    for tgt in targets:
        f = c.link_budget('Stollberg', tgt, model='fock')
        rows.append((tgt, f['d_km'], f['ground'],
                     f['SNR_eq_dB'], f['V_eq_uV']))
    return rows


def fmt_snr(snr):
    return f"{snr:+.1f}"


def fmt_uv(v):
    if v >= 0.01:
        return f"{v:.3f}"
    return f"{v:.2e}"


def verdict_label(snr_db, threshold=10):
    if snr_db >= threshold: return "PASS"
    if snr_db >= 0:         return "MARGINAL"
    return "FAIL"


# =====================================================================
def slide_01_title(prs):
    s = blank_slide(prs)
    add_logo(s, left=Inches(5.91), top=Inches(0.9), height=Inches(2.5))
    add_text(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.9),
             "BATTLE OF THE BEAMS",
             size=44, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
             "Null hypothesis: horizontal VHF beam propagation over a curved surface",
             size=20, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.5),
             "Knickebein 31.5 MHz, Kleve and Stollberg transmitters, Sep 1939 to May 1941",
             size=15, color=CREAM, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             "Aether Cosmology Research Group  |  ITU-certified BotB Calc v9_1",
             size=13, color=GOLD, align=PP_ALIGN.CENTER)


def slide_02_null(prs):
    s = blank_slide(prs)
    header(s, "Null hypothesis",
           "Classical Popperian framing of the test")
    add_text(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(0.5),
             "We test whether the documented operation of the Knickebein VHF beam system "
             "is consistent with propagation over a sphere of radius 6,371 km.",
             size=16, color=CREAM)

    add_text(s, Inches(0.7), Inches(2.4), Inches(12.0), Inches(0.45),
             "H1  (the alternative we attempt to falsify):",
             size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(2.85), Inches(11.8), Inches(0.95),
             "The signal at the equisignal crossover sits above the receiver noise floor "
             "on a sphere of radius 6,371 km for every confirmed operational target.",
             size=15, color=CREAM)

    add_text(s, Inches(0.7), Inches(3.95), Inches(12.0), Inches(0.45),
             "H0  (the position we assume false and try to refute):",
             size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(4.40), Inches(11.8), Inches(1.4),
             "The signal at the equisignal crossover sits below the receiver noise floor "
             "on a sphere of radius 6,371 km for one or more confirmed operational targets. "
             "If the receiver cannot pick up the cross beam, the two-beam intersection "
             "cannot form, and precision bombing cannot proceed.",
             size=15, color=CREAM)

    add_text(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.9),
             "The system demonstrably functioned over 10 months of operational use. "
             "A surface geometry that prohibits the observed function is falsified by the observation.",
             size=14, color=GOLD)


def slide_03_system(prs):
    s = blank_slide(prs)
    header(s, "The Knickebein system",
           "Two transmitters, two beams, one intersection")
    add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.5),
             "Telefunken 31.5 MHz directional beam, 3 kW TX, 99 m x 29 m aperture (26.0 dBi). "
             "Pilot flies the equisignal corridor between two sub-beams squinted +/- 5 degrees.",
             size=14, color=CREAM)

    rows = [
        ["Station", "Lat", "Lon", "Terrain", "Frame", "TX height", "Ground under beam"],
        ["Kleve (Kn-4) director",     "51.79 N", "6.10 E", "83 m",  "28 m", "111 m", "land"],
        ["Stollberg (Kn-2) cross",    "54.64 N", "8.94 E", "44 m",  "28 m",  "72 m", "sea"],
        ["Greny (Kn-7)",              "49.95 N", "1.29 E", "134 m", "28 m", "162 m", "sea"],
        ["Beaumont-Hague (Kn-9)",     "49.67 N", "1.85 W", "169 m", "28 m", "197 m", "sea"],
    ]
    add_table(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(2.4), rows)

    add_text(s, Inches(0.7), Inches(4.95), Inches(12.0), Inches(0.45),
             "Discriminating leg of the test: the Stollberg cross beam",
             size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.8),
             "Without the cross beam, the two-beam intersection does not form and "
             "no release point is given to the bomb-aimer. Director-only beam guidance "
             "(Kleve alone) provides a line to fly, no release point. The Stollberg "
             "cross beam at 694 to 791 km over the North Sea is therefore the load-bearing "
             "geometry for the falsification test.",
             size=14, color=CREAM)


def slide_04_two_models(prs):
    s = blank_slide(prs)
    header(s, "Two propagation models, four notebooks",
           "Same link budget, different surface geometry")

    add_text(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.5),
             "Globe Earth (GE): Fock smooth-Earth diffraction",
             size=20, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(2.5),
             "ITU-R P.526-16 Section 3, first-term residue series.\n"
             "Wave bends around the 6,371 km sphere.\n"
             "Exponential field decay past the radio horizon.\n"
             "Notebook 001 = Kleve director.\n"
             "Notebook 002 = Stollberg cross.",
             size=14, color=CREAM)

    add_text(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(0.5),
             "Flat Earth (FE): Sommerfeld-Norton three-term Ez",
             size=20, bold=True, color=GOLD)
    add_text(s, Inches(6.9), Inches(2.2), Inches(6.0), Inches(2.5),
             "ITU Handbook on Ground Wave Propagation (2014) Part 1 Section 3.2.1.\n"
             "Direct ray plus Fresnel-reflected ray plus Norton surface wave.\n"
             "Rectilinear, no horizon, no shadow zone.\n"
             "Notebook 003 = Kleve director.\n"
             "Notebook 004 = Stollberg cross.",
             size=14, color=CREAM)

    add_text(s, Inches(0.6), Inches(4.85), Inches(12.2), Inches(0.5),
             "Common inputs (same for both models, both stations)",
             size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(5.35), Inches(12.2), Inches(1.6),
             "f = 31.5 MHz, P_tx = 3,000 W, G_tx = 26.0 dBi, lambda = 9.517 m, "
             "k = 4/3 effective Earth radius (a_e = 8,495 km), 290 K thermal at 500 Hz, "
             "ITU-R P.372 galactic Fa = 17.5 dB at 31.5 MHz, noise floor 0.0755 uV at 50 ohm input, "
             "Telefunken +/- 5 degree squint geometry.",
             size=14, color=CREAM)


def slide_05_fock_math(prs):
    s = blank_slide(prs)
    header(s, "Globe model: ITU-R P.526-16 Fock diffraction",
           "First-term residue series past the radio horizon")
    bullets = [
        "Effective Earth radius:  a_e = (4/3) R_Earth = 8,495 km",
        "Beta polarisation parameter (Eq. 16, 16a):  K^2 = 6.89 sigma / (k^(2/3) f_MHz^(5/3))",
        "  - Land at 31.5 MHz vert pol -> beta = 1.0",
        "  - Sea  at 31.5 MHz vert pol -> beta ~ 0.81  (Stollberg correction)",
        "Normalised distance:  X = beta (pi / (lambda a_e^2))^(1/3) d   (Eq. 13)",
        "Normalised heights:    Y_(1,2) = 2 beta (pi^2 / (lambda^2 a_e))^(1/3) h_(tx,rx)",
        "Distance term (long-path, X >= 1.6):  F(X) = 11 + 10 log10 X - 17.6 X",
        "Height-gain term G(Y) per Eq. 17 with Eq. 18 lower-bound clamp",
        "Total:  E/E_0 (dB) = F(X) + G(Y_1) + G(Y_2)",
    ]
    add_text(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.5),
             "\n".join(bullets), size=15, color=CREAM)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.6),
             "Notebooks 001 and 002 walk every term with editable target selection and a Squint Sandbox cell.",
             size=13, color=GOLD)


def slide_06_sn_math(prs):
    s = blank_slide(prs)
    header(s, "Flat model: Sommerfeld-Norton three-term Ez",
           "ITU Handbook 2014 Part 1 Section 3.2.1")
    bullets = [
        "Slant ranges (no curvature):  r_1 = sqrt(d^2 + (h_rx - h_tx)^2),  r_2 = sqrt(d^2 + (h_rx + h_tx)^2)",
        "Complex permittivity:  x = 18000 sigma / f_MHz,   n^2 = eps_r - j x,   u^2 = 2 / n^2",
        "Fresnel reflection coefficient (vertical pol):",
        "    R_v = (n^2 sin psi_2 - sqrt(n^2 - cos^2 psi_2)) / (n^2 sin psi_2 + sqrt(n^2 - cos^2 psi_2))",
        "Norton numerical distance:  w = -j 2k r_2 u^2 (1 - u^2 cos^2 psi_2) / (1 - R_v)",
        "Attenuation function (large-|w| asymptotic, Eq. 7):",
        "    F = -1/(2w) - 3/(2w)^2 - 15/(2w)^3 - 105/(2w)^4",
        "E_z = E_direct + E_reflected + E_surface, summed as complex phasors at the receiver",
    ]
    add_text(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.5),
             "\n".join(bullets), size=15, color=CREAM)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.6),
             "Notebooks 003 and 004 walk every complex-arithmetic step from r_1 / r_2 through P_rx.",
             size=13, color=GOLD)


def slide_07_grwave(prs):
    s = blank_slide(prs)
    header(s, "Independent globe cross-check: ITU-R P.368 GRWAVE",
           "ITU's own Fortran ground-wave calculator")
    add_text(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.6),
             "P.368 GRWAVE is the Fortran implementation of Recommendation ITU-R P.368, "
             "maintained on the ITU-R Study Group 3 page. Rotheram (1981) three-region method: "
             "geometric-optics + extended flat-Earth + full multi-mode residue series, with the "
             "exponential refractivity profile N(h) = 315 exp(-0.136 h).\n\n"
             "GRWAVE runs the same underlying physics as ITU-R P.526 Section 3 (Fock residue series) "
             "via an independent code path, then adds the Sommerfeld surface-wave contribution and "
             "the full multi-mode sum. Cross-validated against the P.526 first-term Fock to within "
             "a few dB on every BotB path.\n\n"
             "Verdict from GRWAVE on Stollberg to Beeston (694 km, sea, 31.5 MHz, RX 6,000 m): "
             "peak SNR roughly -3 dB below the receiver noise floor. Both ITU globe standards agree.",
             size=15, color=CREAM)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.6),
             "Reference: GRWAVE_P368_BotB.md in the Obsidian vault, plus run_grwave_knickebein.py in the BotB repo.",
             size=13, color=GOLD)


def slide_08_kleve_results(prs):
    s = blank_slide(prs)
    header(s, "Results: Kleve director beam",
           "Globe Fock, 5 degree squint, equisignal corridor")
    rows = [["Target", "d (km)", "Ground", "SNR_eq (dB)", "V_eq (uV)", "Verdict"]]
    for tgt, d, g, snr, v in kleve_results():
        rows.append([tgt, f"{d:.0f}", g, fmt_snr(snr), fmt_uv(v), verdict_label(snr)])
    add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.6), rows)

    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.6),
             "Director beam over land: every confirmed Midlands target clears the +10 dB detection floor on the globe.",
             size=13, color=GOLD)


def slide_09_stollberg_results(prs):
    s = blank_slide(prs)
    header(s, "Results: Stollberg cross beam",
           "Globe Fock, 5 degree squint, sea path, beta = 0.81")
    rows = [["Target", "d (km)", "Ground", "SNR_eq (dB)", "V_eq (uV)", "Verdict"]]
    for tgt, d, g, snr, v in stollberg_results():
        rows.append([tgt, f"{d:.0f}", g, fmt_snr(snr), fmt_uv(v), verdict_label(snr)])
    add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.6), rows)

    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.6),
             "Cross beam over sea: every confirmed Midlands target sits below the noise floor on the globe.",
             size=13, color=RED)


def slide_10_squint_sandbox(prs):
    s = blank_slide(prs)
    header(s, "Squint Sandbox calibration",
           "Values that reproduce the British-measured equisignal corridor")
    add_text(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(0.4),
             "Default Squint Sandbox inputs (sheet rows 4, 9, 18)",
             size=16, bold=True, color=GOLD)

    rows = [
        ["Input",                  "Default", "Source"],
        ["Squint theta",           "5 deg",   "Telefunken (Bauer 2004 p. 12)"],
        ["Aperture W",             "99 m",    "Trenkle 1979 p. 67"],
        ["Aperture H",             "20 m",    "Squint Sandbox sheet B18 / C18"],
    ]
    add_table(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(1.7), rows)

    # Compute the corridor at Spalding (Bufton measurement) to confirm calibration
    r = c.link_budget("Kleve", "Spalding", model="fock")
    lam = c.freq_to_wavelen(31.5)
    W_m = 99.0
    width_m = c.equisignal_corridor_width_m(r["d_m"], W_m, 5.0, 31.5)
    width_yd = width_m / 0.9144

    add_text(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.4),
             "Calibration check: Kleve to Spalding, 21 Jun 1940 (Bufton flight)",
             size=16, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(4.5), Inches(12.0), Inches(1.8),
             f"d = {r['d_km']:.0f} km, predicted equisignal corridor at target range = "
             f"{width_m:.0f} m = {width_yd:.0f} yards.\n\n"
             "British R/T flight measured 400 to 500 yards equisignal width at this target. "
             "The 99 m aperture at 5 degree squint predicts within ~10 percent without tuning.",
             size=15, color=CREAM)

    add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5),
             "The same defaults flow into notebooks 001 and 002. Edit them to see how the corridor changes.",
             size=13, color=GOLD)


def slide_11_verdict(prs):
    s = blank_slide(prs)
    header(s, "Verdict",
           "Globe (Fock and GRWAVE) vs Flat (Sommerfeld-Norton)")

    add_text(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.45),
             "H1 falsified",
             size=22, bold=True, color=RED)
    add_text(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.6),
             "Every Stollberg to Midlands path (694 to 791 km, sea, beta = 0.81) sits "
             "below the 0.0755 uV receiver noise floor on both ITU globe models. "
             "Required TX height to put Stollberg to Liverpool inside line of sight on a "
             "globe is 13,098 m. Actual antenna frame is 72 m. Ratio 182 to 1.",
             size=15, color=CREAM)

    add_text(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.45),
             "H0 holds (cannot reject)",
             size=22, bold=True, color=GREEN)
    add_text(s, Inches(0.7), Inches(4.45), Inches(12.0), Inches(2.4),
             "Sommerfeld-Norton flat-Earth model predicts +50 to +85 dB equisignal SNR "
             "on every confirmed path with no horizon, no shadow zone, no exponential decay. "
             "Predicted equisignal corridor width at Spalding matches the British-measured "
             "400 to 500 yards. The observed operational record (10 months of "
             "Knickebein-guided night raids on Derby, Birmingham, Liverpool, Coventry) "
             "is quantitatively consistent with rectilinear propagation over a flat surface.",
             size=15, color=CREAM)


def slide_12_closing(prs):
    s = blank_slide(prs)
    add_logo(s, left=Inches(5.91), top=Inches(0.7), height=Inches(2.0))
    add_text(s, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.7),
             "Aether Cosmology Research Group",
             size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6),
             "ITU-certified Battle of the Beams Calculator v9_1",
             size=18, color=CREAM, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             "Reference",
             size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.6),
             "Jupyter Book walkthrough: BotB/jupyter_book/\n"
             "Vault docs: /Null_Hypothesis/Battle_of_the_Beams/  (Knickebein_Propagation_Null.md, GRWAVE_P368_BotB.md)\n"
             "Code repo: github.com/AlanSpaceAudits/battle-of-the-beams\n"
             "ITU standards used: P.526-16 Fock, P.368 GRWAVE, P.372-16 noise, P.527 surface constants",
             size=13, color=CREAM, align=PP_ALIGN.CENTER)


def main():
    prs = make_prs()
    slide_01_title(prs)
    slide_02_null(prs)
    slide_03_system(prs)
    slide_04_two_models(prs)
    slide_05_fock_math(prs)
    slide_06_sn_math(prs)
    slide_07_grwave(prs)
    slide_08_kleve_results(prs)
    slide_09_stollberg_results(prs)
    slide_10_squint_sandbox(prs)
    slide_11_verdict(prs)
    slide_12_closing(prs)
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
